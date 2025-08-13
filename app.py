import pandas as pd
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
import os

class ExcelToJsonPresentation:
    def __init__(self, excel_file_path):
        """
        初始化類別
        
        Args:
            excel_file_path (str): Excel文件路徑
        """
        self.excel_file_path = excel_file_path
        self.json_data = None
        self.filtered_total = 0
        
    def _set_font_style(self, text_frame, font_name="微軟正黑體", font_size=Pt(12)):
        """
        設定文字框架的字體樣式
        
        Args:
            text_frame: 文字框架
            font_name (str): 字體名稱
            font_size: 字體大小
        """
        for paragraph in text_frame.paragraphs:
            # 設定段落字體
            paragraph.font.name = font_name
            paragraph.font.size = font_size
            
            # 設定runs字體
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = font_size
            
            # 如果段落沒有runs，添加一個並設定字體
            if not paragraph.runs:
                run = paragraph.add_run()
                run.font.name = font_name
                run.font.size = font_size
    
    def _set_paragraph_font(self, paragraph, font_name="微軟正黑體", font_size=Pt(12)):
        """
        設定段落的字體樣式
        
        Args:
            paragraph: 段落對象
            font_name (str): 字體名稱
            font_size: 字體大小
        """
        if paragraph.runs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = font_size
        else:
            # 如果段落沒有runs，創建一個
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            run.font.name = font_name
            run.font.size = font_size
        
    def read_excel_to_json(self, sheet_name=None):
        """
        讀取Excel文件並轉換為JSON格式
        
        Args:
            sheet_name (str, optional): 工作表名稱，預設為None（讀取第一個工作表）
            
        Returns:
            dict: 轉換後的JSON數據
        """
        try:
            # 讀取Excel文件，使用openpyxl引擎以更好地處理資料驗證
            if sheet_name:
                df = pd.read_excel(self.excel_file_path, sheet_name=sheet_name, engine='openpyxl')
            else:
                df = pd.read_excel(self.excel_file_path, engine='openpyxl')
            
            # 使用openpyxl直接讀取Excel以獲取資料驗證欄位的實際值
            from openpyxl import load_workbook
            wb = load_workbook(self.excel_file_path, data_only=True)
            if sheet_name:
                ws = wb[sheet_name]
            else:
                ws = wb.active
            
            # 重新構建DataFrame，確保讀取到實際的儲存格值
            data_rows = []
            headers = [cell.value if cell.value is not None else f'Unnamed_{i}' for i, cell in enumerate(ws[1])]  # 第一行作為標題，處理None值
            
            for row in ws.iter_rows(min_row=2, values_only=True):
                if any(cell is not None for cell in row):  # 跳過完全空白的行
                    data_rows.append(row)
            
            # 創建新的DataFrame
            df = pd.DataFrame(data_rows, columns=headers)
            
            # 清理數據：移除空值行
            df = df.dropna(how='all')
            
            # 處理NaN值和None值：將其替換為空字串
            df = df.fillna('')
            df = df.replace({None: ''})
            
            # 按類型欄位分類議題
            categorized_data = {}
            if '類型' in df.columns:
                # 按類型分組
                grouped = df.groupby('類型')
                for category, group in grouped:
                    if "宣導" in category:
                        continue
                    else:
                        categorized_data[category] = group.to_dict('records')
                
                # 計算各類型的統計資訊
                category_stats = {}
                for category, items in categorized_data.items():
                    if "宣導" in category:
                        continue
                    else:
                        category_stats[category] = len(items)
            else:
                # 如果沒有類型欄位，將所有資料放在 "未分類" 中
                categorized_data['未分類'] = df.to_dict('records')
                category_stats = {'未分類': len(df)}
            
            # 轉換為JSON格式
            # 排除含「宣導」之類別於統計清單
            filtered_categories_for_meta = [c for c in categorized_data.keys() if '宣導' not in c]
            json_data = {
                'metadata': {
                    'file_name': os.path.basename(self.excel_file_path),
                    'total_rows': len(df),
                    'columns': list(df.columns),
                    'sheet_name': sheet_name if sheet_name else 'Sheet1',
                    'categories': filtered_categories_for_meta,
                    'category_stats': category_stats
                },
                'data_by_category': categorized_data,
                'data': df.to_dict('records')  # 保留原始格式以便向後相容
            }
            
            self.json_data = json_data
            return json_data
            
        except Exception as e:
            print(f"讀取Excel文件時發生錯誤: {str(e)}")
            return None
    
    def save_json(self, output_path):
        """
        將JSON數據保存到文件
        
        Args:
            output_path (str): 輸出JSON文件路徑
        """
        if self.json_data:
            try:
                with open(output_path, 'w', encoding='utf-8') as f:
                    json.dump(self.json_data, f, ensure_ascii=False, indent=2)
                print(f"JSON文件已保存至: {output_path}")
                return True
            except Exception as e:
                print(f"保存JSON文件時發生錯誤: {str(e)}")
                return False
        else:
            print("沒有JSON數據可保存，請先執行 read_excel_to_json()")
            return False
    
    def create_presentation(self, output_pptx_path, items_per_page=3):
        """
        根據JSON數據創建PowerPoint簡報
        
        Args:
            output_pptx_path (str): 輸出PPT文件路徑
            items_per_page (int): 每頁顯示的議題數量，預設為3
        """
        if not self.json_data:
            print("沒有JSON數據可用於創建簡報，請先執行 read_excel_to_json()")
            return False
        
        try:
            # 創建新的簡報
            prs = Presentation()
            # 計算排除「宣導」後的總議題數
            try:
                if 'data_by_category' in self.json_data and isinstance(self.json_data['data_by_category'], dict):
                    for category, items in self.json_data['data_by_category'].items():
                        if '宣導' in category:
                            continue
                        self.filtered_total += len(items or [])
                else:
                    for item in self.json_data.get('data', []) or []:
                        category = str(item.get('類型', '') or '')
                        if '宣導' in category:
                            continue
                        self.filtered_total += 1
            except Exception:
                self.filtered_total = self.json_data.get('metadata', {}).get('total_rows', 0)
            
            # 設置簡報標題頁
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = f"系統交流會議題討論 - {self.json_data['metadata']['file_name'].replace('.xlsx', '')}"
            subtitle.text = f"共 {self.filtered_total} 筆議題"
            
            # 設定標題頁字體
            self._set_font_style(title.text_frame, font_size=Pt(30))
            self._set_font_style(subtitle.text_frame, font_size=Pt(20))
            bullet_slide_layout = prs.slide_layouts[1]
            """
            # 目錄頁（第二頁）
            categories = self.json_data['metadata'].get('categories', [])
            if (not categories) and ('data_by_category' in self.json_data):
                categories = list(self.json_data['data_by_category'].keys())
            self._add_toc_slide(prs, categories, bullet_slide_layout)
            """
            # 議題動機頁（第三頁）
            self._add_motivation_slide(prs, bullet_slide_layout)

            # 添加概覽頁
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "議題概覽"
            tf = content.text_frame
            tf.text = f"總議題數: {self.filtered_total}"
            
            # 顯示各類型統計
            if 'category_stats' in self.json_data['metadata']:
                p = tf.add_paragraph()
                p.text = "各類型議題統計:"
                for category, count in self.json_data['metadata']['category_stats'].items():
                    if "宣導" in category:
                        continue
                    else:
                        p = tf.add_paragraph()
                        p.text = f"  {category}: {count} 筆"
                        p.level = 1
            
            # 設定概覽頁字體
            self._set_font_style(title.text_frame, font_size=Pt(28))
            self._set_font_style(tf, font_size=Pt(20))
            
            # 按類型分組顯示議題
            if 'data_by_category' in self.json_data:
                for category, items in self.json_data['data_by_category'].items():
                    if "宣導" in category:
                        continue
                    else:
                        self._add_category_slides(prs, category, items, items_per_page, bullet_slide_layout)
            else:
                # 如果沒有分類數據，使用原始數據
                self._add_category_slides(prs, "所有議題", self.json_data['data'], items_per_page, bullet_slide_layout)
            
            # 保存簡報（若原檔案被占用，改另存新檔）
            target_path = output_pptx_path
            try:
                prs.save(target_path)
            except Exception as save_e:
                if 'Permission denied' in str(save_e) or isinstance(save_e, PermissionError):
                    base, ext = os.path.splitext(output_pptx_path)
                    alt_path = f"{base}_另存{ext}"
                    counter = 1
                    while os.path.exists(alt_path):
                        alt_path = f"{base}_另存({counter}){ext}"
                        counter += 1
                    prs.save(alt_path)
                    target_path = alt_path
                    print(f"原檔案被占用，已另存為: {target_path}")
                else:
                    raise
            print(f"簡報已保存至: {target_path}")
            # 嘗試在 Windows 上透過 PowerPoint COM 套用佈景主題（預設為第一個可用佈景主題/Office Theme）
            try:
                self._apply_theme_windows(target_path)
            except Exception as theme_e:
                print(f"套用佈景主題時發生警告：{theme_e}")
            return True
            
        except Exception as e:
            print(f"創建簡報時發生錯誤: {str(e)}")
            return False

    def _add_toc_slide(self, prs, categories, bullet_slide_layout):
        """
        添加目錄頁，列出所有類型
        
        Args:
            prs: PowerPoint簡報對象
            categories (list[str]): 類型清單
            bullet_slide_layout: 簡報佈局
        """
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "目錄"
        tf = content.text_frame
        tf.clear()

        # 使用提供的 categories，若空則從 data_by_category 推得
        if (not categories) and ('data_by_category' in self.json_data):
            categories = list(self.json_data['data_by_category'].keys())

        # 過濾不需要的分類（例如：包含「宣導」）
        filtered = [c for c in categories if '宣導' not in c]

        # 期望排序（若存在則依此順序）
        preferred_order = [
            '系統問題',
            '系統優化',
            '資料一致性問題',
            '跨部門協作與流程建議',
            '資安與基礎設施政策',
            '權限設定與管理議題',
            '資料公告',
        ]
        ordered = [c for c in preferred_order if c in filtered]
        ordered += [c for c in filtered if c not in ordered]

        # 附帶數量統計
        stats = (self.json_data.get('metadata') or {}).get('category_stats') or {}

        # 寫入目錄內容
        if ordered:
            for idx, category in enumerate(ordered, 1):
                count = stats.get(category)
                label = f"{category} ({count} 項)" if isinstance(count, int) else category
                if idx == 1:
                    tf.text = label
                else:
                    p = tf.add_paragraph()
                    p.text = label
        else:
            tf.text = "尚無分類資料"
        self._set_font_style(title.text_frame, font_size=Pt(28))
        self._set_font_style(tf, font_size=Pt(20))

    def _add_motivation_slide(self, prs, bullet_slide_layout):
        """
        添加「議題動機」頁：拆分為三頁 PART 1 / PART 2 / PART 3
        """
        # PART 1
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "議題動機 - PART 1. 現況與挑戰"
        tf = content.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.text = "現況與挑戰"
        p = tf.add_paragraph()
        p.text = f"隨著業務的持續推展，我們在日常作業流程與資訊系統上累積了 {self.filtered_total} 項待解決的議題。這些挑戰已對部分業務的運作效率、資料準確性與跨部門協作造成影響。"
        p.level = 1
        p = tf.add_paragraph()
        p.text = "為了提升整體營運效能，我們需要集中資源，共同檢視並解決這些關鍵問題。"
        p.level = 1
        self._set_font_style(title.text_frame, font_size=Pt(28))
        self._set_font_style(tf, font_size=Pt(20))

        # PART 2
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "議題動機 - PART 2. 會議目標"
        tf = content.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.text = "會議目標"
        p = tf.add_paragraph()
        p.text = "凝聚共識： 對各項議題的癥結點與解決方案達成共識。"
        p.level = 1
        p = tf.add_paragraph()
        p.text = "確立方向： 明確後續系統優化與流程改善的優先次序。"
        p.level = 1
        p = tf.add_paragraph()
        p.text = "分配任務： 釐清各項任務的負責單位與預計完成時程。"
        p.level = 1
        self._set_font_style(title.text_frame, font_size=Pt(28))
        self._set_font_style(tf, font_size=Pt(20))

        # PART 3
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "議題分類說明"
        tf = content.text_frame
        tf.clear()
        tf.word_wrap = True

        categories_info = [
            ("系統問題", "指系統現有的功能出現錯誤(Bug)、運算失準或流程中斷，導致使用者無法完成特定作業，或產出不正確的結果。這些是需要優先修復的具體技術缺陷。"),
            ("系統優化", "指系統功能雖可運作，但在操作流程、使用者體驗或執行效率上，仍有改善空間。這類建議是為了讓系統更易用、更高效，例如新增便利功能或調整操作邏輯。"),
            ("資料一致性問題", "指不同系統模組或作業環節中，對於相同資料的定義、格式或驗證標準不一，導致資料混亂、轉換錯誤或比對困難。"),
            ("跨部門協作與流程建議", "聚焦於資訊單位與各業務單位之間的溝通、權責劃分與作業流程。旨在建立標準化的需求提出、測試回饋及上線溝通機制，以提升協作效率。"),
            ("資安與基礎設施政策", "涉及影響全公司的資訊安全政策、硬體設備及基礎軟體環境的通用性規範。這類議題屬於原則性、全面性的管理政策。"),
            ("權限設定與管理議題", "探討各系統的權限劃分與管理機制。核心問題是權限設定未與使用者職務需求精準對應，存在過高或不足的狀況，衍生管理不便與潛在資安風險。")
        ]
        # 以每頁 3 個分類進行分頁
        chunks = [categories_info[i:i+3] for i in range(0, len(categories_info), 3)]
        for idx, chunk in enumerate(chunks):
            if idx == 0:
                current_title = title
                current_tf = tf
            else:
                slide = prs.slides.add_slide(bullet_slide_layout)
                current_title = slide.shapes.title
                content = slide.placeholders[1]
                current_title.text = "議題分類說明"
                current_tf = content.text_frame
                current_tf.clear()
                current_tf.word_wrap = True

            # 寫入當頁的分類
            first_title, first_desc = chunk[0]
            current_tf.text = first_title
            p = current_tf.add_paragraph()
            p.text = first_desc
            p.level = 1
            for cat_title, cat_desc in chunk[1:]:
                p = current_tf.add_paragraph()
                p.text = cat_title
                p.level = 0
                p = current_tf.add_paragraph()
                p.text = cat_desc
                p.level = 1

            self._set_font_style(current_title.text_frame, font_size=Pt(28))
            self._set_font_style(current_tf, font_size=Pt(20))

    def _add_category_slides(self, prs, category, items, items_per_page, bullet_slide_layout, max_content_length=200):
        """
        為特定類型的議題添加簡報頁面
        
        Args:
            prs: PowerPoint簡報對象
            category (str): 議題類型
            items (list): 該類型的議題列表
            items_per_page (int): 每頁顯示的議題數量
            bullet_slide_layout: 簡報佈局
            max_content_length (int): 每個議題內容的最大字數，超過則獨立成頁
        """
        if not items:
            return
        
        # 依據內容長度拆頁：超過 max_content_length 的議題單獨成頁
        slides_batches = []  # List[List[Tuple[item, seq_number]]]
        buffer_batch = []
        seq_number = 1
        for original_item in items:
            item_content = original_item.get('內容', '') or ''
            if len(item_content) > max_content_length:
                if buffer_batch:
                    slides_batches.append(buffer_batch)
                    buffer_batch = []
                slides_batches.append([(original_item, seq_number)])
            else:
                buffer_batch.append((original_item, seq_number))
                if len(buffer_batch) >= items_per_page:
                    slides_batches.append(buffer_batch)
                    buffer_batch = []
            seq_number += 1
        if buffer_batch:
            slides_batches.append(buffer_batch)

        total_pages = len(slides_batches)
        for page_num, page_items in enumerate(slides_batches, start=1):
            
            # 創建新頁面
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            # 設置標題
            if total_pages > 1:
                title.text = f"{category} ({page_num}/{total_pages})"
            else:
                title.text = category
            
            # 清空內容框架
            tf = content.text_frame
            tf.clear()

            # 設定標題字體
            self._set_font_style(title.text_frame, font_size=Pt(24))
            
            # 添加議題內容
            long_paragraphs = []  # 需降字級的內容段落
            for i, (item, seq_no) in enumerate(page_items):
                if '原規劃刪除所有車型時不可進行規格構造變更，但實務上卻可以接受此類申請案' in item.get('內容', ''):
                    print(item)
                if '紙本線上平台統一編號' in item.get('內容', ''):
                    print(item)
                # 議題標題（使用全域序號）
                if i == 0:
                    tf.text = f"{seq_no}. {item.get('內容', '無內容')}"
                    content_paragraph = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                    p.text = f"{seq_no}. {item.get('內容', '無內容')}"
                    content_paragraph = p

                # 若內容字數超過 300，將該段內容字級調為 18
                if len(item.get('內容', '') or '') > 300:
                    long_paragraphs.append(content_paragraph)
                
                # 添加詳細資訊
                details = []
                if '涉及處別' in item and item['涉及處別']:
                    details.append(f"涉及處別: {item['涉及處別']}")
                if '涉及部門' in item and item['涉及部門']:
                    details.append(f"涉及部門: {item['涉及部門']}")
                if '系統別' in item and item['系統別']:
                    details.append(f"系統別: {item['系統別']}")
                
                if details:
                    p = tf.add_paragraph()
                    p.text = " | ".join(details)
                    p.level = 1  # 設定為第1層級，縮排顯示
                
                # 添加分隔段落（除了最後一個項目）
                if i < len(page_items) - 1:
                    p = tf.add_paragraph()
                    p.text = ""
            
            # 設定內容字體 - 較小的字體以容納更多內容
            self._set_font_style(tf, font_size=Pt(20))
            # 針對超長內容的項目，將字體調小至 18
            for para in long_paragraphs:
                self._set_paragraph_font(para, font_size=Pt(18))

    def _apply_theme_windows(self, pptx_path, preferred_theme_name="Office"):
        """
        使用 Windows PowerPoint COM 在既有檔案上套用佈景主題。
        會尋找系統中常見的佈景主題資料夾，優先套用檔名包含 preferred_theme_name 的 .thmx，
        找不到則套用第一個可用的 .thmx。
        """
        try:
            import win32com.client as win32
        except Exception:
            print("未安裝 pywin32 或當前環境不支援 COM，自動套用佈景主題略過。")
            return

        # 先將路徑正規化為絕對路徑並確認檔案存在
        pptx_abs_path = os.path.abspath(pptx_path)
        if not os.path.isfile(pptx_abs_path):
            print(f"找不到簡報檔案，略過自動套用佈景主題：{pptx_abs_path}")
            return

        theme_dirs = []

        # 1) 專案根目錄指定主題（優先）
        project_theme = os.path.abspath(os.path.join(os.path.dirname(__file__), '佈景主題.thmx'))
        if os.path.isfile(project_theme):
            selected = project_theme
        else:
            selected = None

        # 使用者自訂主題
        user_profile = os.environ.get('USERPROFILE') or os.path.expanduser('~')
        if user_profile:
            theme_dirs.append(os.path.join(user_profile, 'AppData', 'Roaming', 'Microsoft', 'Templates', 'Document Themes'))

        # Office 常見安裝主題路徑
        pf = os.environ.get('ProgramFiles', r'C:\Program Files')
        pf86 = os.environ.get('ProgramFiles(x86)', r'C:\Program Files (x86)')
        theme_dirs.append(os.path.join(pf, 'Microsoft Office', 'root', 'Document Themes 16', 'Theme Gallery'))
        theme_dirs.append(os.path.join(pf86, 'Microsoft Office', 'root', 'Document Themes 16', 'Theme Gallery'))
        for ver in ['15', '14']:
            theme_dirs.append(os.path.join(pf, 'Microsoft Office', f'Document Themes {ver}', 'Theme Gallery'))
            theme_dirs.append(os.path.join(pf86, 'Microsoft Office', f'Document Themes {ver}', 'Theme Gallery'))

        candidate_thmx = []
        if selected is None:
            for d in theme_dirs:
                try:
                    if os.path.isdir(d):
                        for name in os.listdir(d):
                            if name.lower().endswith('.thmx'):
                                candidate_thmx.append(os.path.join(d, name))
                except Exception:
                    pass

            if not candidate_thmx:
                print("找不到任何 .thmx 佈景主題，略過自動套用。")
                return

            # 優先選擇包含 preferred_theme_name 的主題，其次取第一個
            if preferred_theme_name:
                for p in candidate_thmx:
                    if preferred_theme_name.lower() in os.path.basename(p).lower():
                        selected = p
                        break
            if not selected:
                selected = candidate_thmx[0]

        try:
            app = win32.Dispatch("PowerPoint.Application")
            # 可視化與警示關閉，避免彈窗干擾
            try:
                app.Visible = 0
                app.DisplayAlerts = 0
            except Exception:
                pass

            # 使用絕對路徑開啟簡報
            presentation = app.Presentations.Open(pptx_abs_path, False, False, False)

            # 優先嘗試 ApplyTheme，不支援時改用 ApplyTemplate
            try:
                presentation.ApplyTheme(selected)
            except Exception:
                presentation.ApplyTemplate(selected)

            presentation.Save()
            presentation.Close()
            app.Quit()
            print(f"已套用佈景主題: {os.path.basename(selected)} -> {selected}")
        except Exception as e:
            print(f"PowerPoint COM 自動套用佈景主題失敗：{e}\n- 簡報路徑: {pptx_abs_path}\n- 主題路徑: {selected}")

def main():
    """主函數示例"""
    # Excel文件路徑
    excel_file = "V:/行政服務部/協同作業/資訊/內部/01.資通系統/待討論議題.xlsx"
    
    # 若找不到 Excel，改用現有的 JSON 作為輸入
    if not os.path.exists(excel_file):
        print(f"找不到Excel文件: {excel_file}，改用 output_data.json 產生簡報。")
        json_path = os.path.abspath('output_data.json')
        if not os.path.isfile(json_path):
            print(f"也找不到 JSON 檔：{json_path}")
            return
        try:
            with open(json_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
        except Exception as e:
            print(f"讀取 JSON 失敗：{e}")
            return
        processor = ExcelToJsonPresentation(excel_file)
        processor.json_data = data
        pptx_output = "系統交流會議題簡報.pptx"
        print("正在根據 JSON 直接創建簡報...")
        processor.create_presentation(pptx_output, items_per_page=2)
        print("處理完成！")
        return

    # 一般流程：讀 Excel -> 存 JSON -> 產生 PPTX
    processor = ExcelToJsonPresentation(excel_file)
    print("正在讀取Excel文件...")
    json_data = processor.read_excel_to_json()
    if not json_data:
        print("Excel讀取失敗！")
        return
    print("Excel讀取成功！")
    print(f"共讀取 {json_data['metadata']['total_rows']} 筆資料")
    print(f"欄位: {', '.join(json_data['metadata']['columns'])}")
    json_output = "output_data.json"
    processor.save_json(json_output)
    pptx_output = "系統交流會議題簡報.pptx"
    print("正在創建簡報...")
    processor.create_presentation(pptx_output, items_per_page=2)
    print("處理完成！")

if __name__ == "__main__":
    main()